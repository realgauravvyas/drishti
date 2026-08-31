"""Populate the ICSIGT 2026 conference poster template with the DRISHTI project.

The template's own frames, section bars, logos and colour scheme are left
untouched — only placeholder text is replaced and figures are dropped into the
figure boxes. Output is a 36 x 48 in print-ready poster.
"""
import os, sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(HERE, "img")
ROOT = os.path.join(HERE, "..", "..")
TEMPLATE = os.path.join(
    ROOT, "1st_International_Conference_on_Sustainable_Innovation_and_Green.pptx")
OUT = os.path.join(HERE, "DRISHTI_Poster_ICSIGT2026_EcoGenesis.pptx")

E = 914400.0
NAVY = RGBColor(0x08, 0x30, 0x70)
GREY = RGBColor(0x6E, 0x6E, 0x6E)
BLACK = RGBColor(0x11, 0x11, 0x11)
RED = RGBColor(0xC0, 0x00, 0x00)
GREEN = RGBColor(0x00, 0xB0, 0x50)
ARIAL = "Arial"


# ------------------------------------------------------------------ helpers
def set_text(shape, blocks, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """blocks: list of (text, size, bold, color) or a plain string."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(blocks, str):
        blocks = [(blocks, 18, False, GREY)]
    # wipe every existing paragraph except the first
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    first = tf.paragraphs[0]
    for r in list(first.runs):
        r._r.getparent().remove(r._r)

    for i, blk in enumerate(blocks):
        text, size, bold, color = blk[0], blk[1], blk[2], blk[3]
        space = blk[4] if len(blk) > 4 else 6
        p = first if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        p.line_spacing = 1.12
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = ARIAL
        run.font.color.rgb = color
    return shape


def blank(shape):
    set_text(shape, [("", 12, False, GREY)])


def fit_picture(slide, png, L, T, W, H, pad=0.14, dpi=200.0):
    """Drop an image centred inside a box, preserving aspect ratio."""
    from PIL import Image
    iw, ih = Image.open(png).size
    aw, ah = W - 2 * pad, H - 2 * pad
    scale = min(aw / (iw / dpi), ah / (ih / dpi))
    w, h = (iw / dpi) * scale, (ih / dpi) * scale
    left = L + (W - w) / 2.0
    top = T + (H - h) / 2.0
    return slide.shapes.add_picture(png, Inches(left), Inches(top),
                                    width=Inches(w), height=Inches(h))


def place(slide, png, L, T, W):
    """Place an image at an exact width, height derived from its aspect."""
    from PIL import Image
    iw, ih = Image.open(png).size
    return slide.shapes.add_picture(png, Inches(L), Inches(T),
                                    width=Inches(W), height=Inches(W * ih / iw))


def move(shape, L=None, T=None, W=None, H=None):
    if L is not None: shape.left = Inches(L)
    if T is not None: shape.top = Inches(T)
    if W is not None: shape.width = Inches(W)
    if H is not None: shape.height = Inches(H)


# =============================================================== CONTENT
TITLE = ("DRISHTI: Resolving the Post-Disaster Information Fog "
         "by Treating Silence as Evidence")
AUTHORS = "Gaurav Vyas and Mandavi Singh"
AFFIL = ("Mehta Family School of Data Science and Artificial Intelligence, "
         "Indian Institute of Technology Guwahati")
EMAIL = "g.vyas@op.iitg.ac.in   ·   s.mandavi@op.iitg.ac.in"

ABSTRACT = [
    ("In the first 24 hours after a climate-driven multi-hazard event, a "
     "District Emergency Operations Centre (EOC) receives thousands of "
     "fragmented, contradictory and unverified ground reports, while severed "
     "roads prevent physical verification. Scarce rescue assets are therefore "
     "allocated to wherever reporting is loudest.", 18, False, GREY),
    ("We identify a systematic failure in this practice: report volume is "
     "non-monotonic in damage. Undamaged settlements generate no reports "
     "because there is nothing to report, and catastrophically damaged "
     "settlements also generate none because casualties and destroyed "
     "telecom infrastructure remove the reporters themselves. Volume-based "
     "triage therefore ranks destroyed villages alongside safe ones.",
     18, False, GREY),
    ("DRISHTI reframes the task from establishing truth to maintaining a "
     "calibrated belief. Reports are geo-resolved to a probability "
     "distribution over settlements, parsed for claims, clustered so that "
     "amplified rumour cannot masquerade as corroboration, and fused by "
     "Dempster-Shafer combination of simple support functions in closed "
     "form. Critically, ABSENCE is converted into positive evidence by "
     "fusing telecom-tower liveness, chatter deficit against a population "
     "baseline, and terrain exposure.", 18, False, GREY),
    ("Evaluated on a 1,102-settlement district built from OpenStreetMap with "
     "hidden ground truth, DRISHTI reaches 6,057 of 10,928 affected people "
     "within 24 h against 1,006 for volume ranking, using identical assets. "
     "For the 492 settlements that emit no reports at all, report-driven "
     "ranking is statistically indistinguishable from chance (AUC 0.500) "
     "whereas DRISHTI attains 0.848.", 18, False, GREY),
    ("Keywords: disaster informatics, evidence fusion, Dempster-Shafer "
     "theory, climate resilience, humanitarian logistics, open geospatial "
     "data", 18, True, NAVY),
    ("", 8, False, GREY),
    ("KEY OUTCOMES", 22, True, NAVY),
    ("•  6.0x more affected people reached with identical assets", 18, False, GREY),
    ("•  AUC 0.848 vs 0.500 on settlements that never reported", 18, False, GREY),
    ("•  Blackout detection at 64% precision vs 34% base rate", 18, False, GREY),
    ("•  ~900 of 3,256 reports suppressed as panic amplification", 18, False, GREY),
    ("•  Runs offline on one laptop  ·  zero licence or cloud cost",
     18, False, GREY),
]

OBJECTIVES = [
    ("1.  Quantify the non-monotonic relationship between damage severity "
     "and citizen-report volume, and its effect on asset allocation.",
     18, False, GREY),
    ("2.  Fuse heterogeneous, contradictory reports into a per-settlement "
     "belief distribution that states its own confidence.", 18, False, GREY),
    ("3.  Convert the ABSENCE of expected reports into positive, calibrated "
     "evidence of catastrophe, separating “quiet because unharmed” from "
     "“quiet because destroyed”.", 18, False, GREY),
    ("4.  Suppress rumour amplification so that repetition is not mistaken "
     "for independent corroboration.", 18, False, GREY),
    ("5.  Produce two distinct operational queues — rescue under confidence, "
     "reconnaissance under uncertainty — routed over a damaged network.",
     18, False, GREY),
    ("6.  Deliver the whole system on free and open data at zero marginal "
     "cost, deployable on existing district hardware.", 18, False, GREY),
]

DISCUSSION = [
    ("Why silence is hard.  An undamaged village and a destroyed village are "
     "both silent, so report volume alone can never separate them. We "
     "resolve this with telecom-tower liveness, an observation about the "
     "PLACE rather than about whether anyone had something to say.",
     18, False, GREY),
    ("Tower loss is not proof.  In our scenario 51% of undamaged settlements "
     "also lose their tower, because the grid is cut district-wide to "
     "prevent electrocution. Tower-down is therefore treated as weak "
     "evidence and never used alone; terrain exposure supplies the prior "
     "that separates structural collapse from a generator that ran dry.",
     18, False, GREY),
    ("Guarding against circular evaluation.  An early version of our "
     "simulator generated damage from the same terrain variables the "
     "engine's prior consumes, yielding an inflated AUC of 0.98. We "
     "introduced a spatially-correlated latent vulnerability field "
     "(building stock, drainage, upstream release) that no free map layer "
     "reveals. The honest figure is 0.857.", 18, False, GREY),
    ("Ablation.  Removing the silence engine costs 0.061 AUC overall and "
     "0.137 on non-reporting settlements, confirming that absence carries "
     "genuine information rather than acting as a tie-breaker.",
     18, False, GREY),
    ("Limitations.  Ground truth is simulated, since no labelled multi-hazard "
     "report corpus exists for an Indian district; population is estimated "
     "from OSM place class where census figures are absent; and the "
     "telecom feed is modelled rather than live. None of these affect the "
     "architecture, only the calibration constants.", 18, False, GREY),
]

CONCLUSIONS = [
    ("1.  Report volume is a hill, not a slope. Ranking by it systematically "
     "de-prioritises the highest-mortality settlements.", 18, False, GREY),
    ("2.  Absence of expected communication, fused with infrastructure "
     "liveness and terrain exposure, is a usable and calibrated signal of "
     "catastrophe.", 18, False, GREY),
    ("3.  On settlements that never report, report-driven triage is "
     "equivalent to a coin flip (AUC 0.500); DRISHTI achieves 0.848.",
     18, False, GREY),
    ("4.  The system reaches 6.0x more affected people within the golden "
     "24 hours using the same vehicles, and wastes fewer deployments.",
     18, False, GREY),
    ("5.  Because it runs offline on free and open data, it is deployable "
     "where connectivity and budget are the binding constraints — the "
     "conditions that define a disaster.", 18, False, GREY),
]

ACK = [
    ("Submitted to ICSIGT 2026 · AVINYA 2026 Eco-Innovate Challenge, "
     "organised by the Prakriti Club, IIT Guwahati.", 17, False, GREY),
    ("We thank the OpenStreetMap contributors, Open-Meteo and the USGS for "
     "the open data this work rests on.", 17, False, GREY),
    ("Team EcoGenesis  ·  MIT Licence", 18, True, NAVY),
]

REFS = [
    ("1.  Shafer, G. A Mathematical Theory of Evidence. Princeton Univ. "
     "Press, 1976.", 15, False, GREY),
    ("2.  Imran, M. et al. Processing social media messages in mass "
     "emergency: a survey. ACM Comput. Surv. 47(4), 2015.", 15, False, GREY),
    ("3.  Wald, D. J. et al. Relationships between peak ground acceleration, "
     "peak ground velocity and Modified Mercalli Intensity. Earthquake "
     "Spectra 15(3), 1999.", 15, False, GREY),
    ("4.  National Disaster Management Authority. National Disaster "
     "Management Plan. Govt. of India, 2019.", 15, False, GREY),
    ("5.  OpenStreetMap contributors. Planet dump, 2026. "
     "openstreetmap.org", 15, False, GREY),
    ("6.  Open-Meteo. Free weather and elevation API, 2026. "
     "open-meteo.com", 15, False, GREY),
    ("7.  USGS. ANSS Comprehensive Earthquake Catalog (ComCat), 2026.",
     15, False, GREY),
]


def main():
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass
    prs = Presentation(TEMPLATE)
    s = prs.slides[0]
    sh = list(s.shapes)

    # ---------------------------------------------------------- header
    # The template stacks the title box on top of the conference-name line,
    # which only works while the title is placeholder-short. Re-flow the whole
    # header so a real two-line title clears it.
    move(sh[1], L=5.00, T=2.62, W=26.00, H=1.30)
    move(sh[2], L=5.00, T=4.02, W=26.00, H=0.58)
    move(sh[3], L=5.00, T=4.62, W=26.00, H=0.60)
    move(sh[4], L=5.00, T=5.32, W=26.00, H=0.55)
    set_text(sh[1], [(TITLE, 40, True, NAVY)], align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    set_text(sh[2], [(AUTHORS, 24, True, BLACK)], align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    set_text(sh[3], [(AFFIL, 20, True, BLACK)], align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    set_text(sh[4], [(EMAIL, 18, True, NAVY)], align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

    # ------------------------------------------------- left column text
    # The two left-column frames are very tall (24.8 in and 14.8 in). Scale
    # the body up and open the leading so each column fills its frame instead
    # of leaving a block of dead white at the bottom.
    def relax(blocks, dsize, space):
        return [(t, sz + dsize, b, c, space) for (t, sz, b, c) in blocks]

    set_text(sh[8], relax(ABSTRACT, 1.5, 13), align=PP_ALIGN.LEFT)
    set_text(sh[12], relax(DISCUSSION, 1.0, 12), align=PP_ALIGN.LEFT)

    # ------------------------------------------------------- objectives
    set_text(sh[20], OBJECTIVES, align=PP_ALIGN.LEFT)

    # ------------------------------------------------ figures into boxes
    blank(sh[16])
    fit_picture(s, os.path.join(IMG, "motivation.png"), 7.03, 7.71, 14.15, 6.67)

    blank(sh[24])
    fit_picture(s, os.path.join(IMG, "method.png"), 7.03, 16.41, 10.35, 6.62)

    blank(sh[28])
    fit_picture(s, os.path.join(IMG, "setup.png"), 18.36, 16.41, 9.15, 6.62)

    blank(sh[32])
    fit_picture(s, os.path.join(IMG, "mechanism.png"), 28.49, 16.41, 7.01, 6.62)

    # RESULT - three panels with captions beneath
    panels = [
        (36, 37, 38, 7.06, "reached.png",
         "Fig. 1  People reached within 24 h under a fixed 40-asset budget."),
        (39, 40, 41, 16.63, "blind.png",
         "Fig. 2  Detection accuracy on the 492 settlements emitting no reports."),
        (42, 43, 44, 26.20, "ablation.png",
         "Fig. 3  Ablation: contribution of each mechanism to district-wide AUC."),
    ]
    for _box, lbl, cap, L, png, caption in panels:
        blank(sh[lbl])
        fit_picture(s, os.path.join(IMG, png), L, 25.17, 9.27, 7.15)
        set_text(sh[cap], [(caption, 15, False, GREY)], align=PP_ALIGN.LEFT)

    # RESULT - two wide panels; move their labels below the frames
    wide = [
        (46, 7.06, "agastmuni.png",
         "Fig. 4  Case study. A real settlement of 19,758 people that emitted "
         "one report where 418 were expected."),
        (48, 21.41, "confusion.png",
         "Fig. 5  Confusion matrix over 1,102 settlements and calibrated "
         "performance after temperature scaling."),
    ]
    for lbl, L, png, caption in wide:
        fit_picture(s, os.path.join(IMG, png), L, 33.72, 14.05, 5.55)
        move(sh[lbl], L=L + 0.15, T=39.32, W=13.76, H=0.62)
        set_text(sh[lbl], [(caption, 15, False, GREY)], align=PP_ALIGN.LEFT)

    # ------------------------------------------------------ bottom row
    set_text(sh[52], CONCLUSIONS, align=PP_ALIGN.LEFT)
    set_text(sh[56], ACK, align=PP_ALIGN.LEFT)
    set_text(sh[60], REFS, align=PP_ALIGN.LEFT)

    # ------------------------------------------------------- QR codes
    # A poster is read standing up, so give people a way to open the live
    # dashboard. These sit in the lower half of the ACKNOWLEDGEMENT frame
    # (19.36-27.41 x, 41.02-47.85 y); the text above is kept short to clear them.
    QW, GAP = 1.72, 0.95
    x0 = 19.36 + (8.05 - (2 * QW + GAP)) / 2.0
    place(s, os.path.join(IMG, "qr_demo.png"), x0, 44.92, QW)
    place(s, os.path.join(IMG, "qr_repo.png"), x0 + QW + GAP, 44.92, QW)
    for cx, head, sub, col in [
            (x0 + QW / 2, "LIVE DASHBOARD",
             "realgauravvyas.github.io\n/drishti", GREEN),
            (x0 + QW + GAP + QW / 2, "SOURCE CODE",
             "github.com/realgauravvyas\n/drishti", NAVY)]:
        tb = s.shapes.add_textbox(Inches(cx - 1.45), Inches(46.74),
                                  Inches(2.90), Inches(0.95))
        set_text(tb, [(head, 14, True, col), (sub, 10, False, GREY)],
                 align=PP_ALIGN.CENTER)

    prs.save(OUT)
    print("saved -> %s" % OUT)
    print("size: %.0f x %.0f in" % (prs.slide_width / E, prs.slide_height / E))


if __name__ == "__main__":
    main()
